import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
import os
from PIL import Image
import logging
from tqdm import tqdm
import multiprocessing
from functools import partial
from io import BytesIO
import json
import base64

def get_fill_color(shape):
    if hasattr(shape, 'fill'):
        fill = shape.fill
        if fill.type == MSO_FILL.SOLID:
            if fill.fore_color.rgb:
                return rgb_to_hex(fill.fore_color.rgb)
        elif fill.type == MSO_FILL.BACKGROUND:
            return 'transparent'
    return None

def extract_slide_data(slide):
    return [
        {
            'shape_type': shape.shape_type,
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
            'has_text_frame': shape.has_text_frame,
            'text_frame': extract_text_frame(shape.text_frame) if shape.has_text_frame else None,
            'image': base64.b64encode(shape.image.blob).decode('utf-8') if shape.shape_type == MSO_SHAPE_TYPE.PICTURE else None,
            'image_format': shape.image.ext if shape.shape_type == MSO_SHAPE_TYPE.PICTURE else None,
            'fill_color': get_fill_color(shape)
        }
        for shape in slide.shapes
    ]

def rgb_to_hex(rgb):
    if rgb is None:
        return None
    return '#{:02x}{:02x}{:02x}'.format(*rgb)

def extract_text_frame(text_frame):
    def rgb_to_hex(rgb):
        if rgb is None or not hasattr(rgb, 'rgb') or rgb.rgb is None:
            return None
        # Convert the RGB tuple to hex
        return '#{:02x}{:02x}{:02x}'.format(*rgb.rgb)

    try:
        return [
            {
                'alignment': paragraph.alignment,
                'is_title': paragraph.level == 0,  # Assuming level 0 is for titles
                'runs': [
                    {
                        'text': run.text,
                        'font_name': run.font.name,
                        'font_size': run.font.size.pt if run.font.size else None,
                        'bold': run.font.bold,
                        'italic': run.font.italic,
                        'color': rgb_to_hex(run.font.color),
                        'hyperlink': run.hyperlink.address if run.hyperlink else None
                    }
                    for run in paragraph.runs
                ]
            }
            for paragraph in text_frame.paragraphs
        ]
    except Exception as e:
        logging.warning(f"Error extracting text frame: {str(e)}")
        return []

def process_slide(slide_data_json, image_dir, slide_index, ignore_images=False):
    slide_data = json.loads(slide_data_json)
    return "".join([process_shape_data(shape_index, shape_data, image_dir, slide_index, ignore_images) 
                    for shape_index, shape_data in enumerate(slide_data)])

def process_shape_data(shape_index, shape_data, image_dir, slide_index, ignore_images=False):
    # Calculate relative positions and sizes
    left_percent = f"{(shape_data['left'] / 9144000) * 100:.2f}%"
    top_percent = f"{(shape_data['top'] / 6858000) * 100:.2f}%"
    width_percent = f"{(shape_data['width'] / 9144000) * 100:.2f}%"
    height_percent = f"{(shape_data['height'] / 6858000) * 100:.2f}%"

    common_style = f"left:{left_percent};top:{top_percent};width:{width_percent};height:{height_percent};"

    if shape_data['shape_type'] == MSO_SHAPE_TYPE.PICTURE:
        # Handle images
        image_filename = f"slide_{slide_index + 1}_image_{shape_index + 1}.png"
        
        if not ignore_images:
            image_bytes = base64.b64decode(shape_data['image'])
            image_path = os.path.join(image_dir, image_filename)
            
            try:
                with Image.open(BytesIO(image_bytes)) as img:
                    img.save(image_path, "PNG")
            except OSError as e:
                logging.warning(f"Failed to process image on slide {slide_index + 1}, shape {shape_index + 1}: {str(e)}")
                return f"<p>[Image processing failed for slide {slide_index + 1}, shape {shape_index + 1}]</p>"
        
        return f"<img src='images/{image_filename}' class='absolute object-contain' style='{common_style}'/>"
    elif shape_data['has_text_frame']:
        # Handle text with styling
        text_content = ""
        for paragraph in shape_data['text_frame']:
            para_style = f"text-align: {paragraph['alignment']};"
            para_content = ""
            for run in paragraph['runs']:
                run_style = f"font-family: {run['font_name']}; "
                run_style += f"font-size: {run['font_size']}pt; " if run['font_size'] else ""
                run_style += "font-weight: bold; " if run['bold'] else ""
                run_style += "font-style: italic; " if run['italic'] else ""
                run_style += f"color: {run['color']}; " if run['color'] else ""
                
                # Handle links
                if run.get('hyperlink'):
                    para_content += f"<a href='{run['hyperlink']}' style='{run_style}'>{run['text']}</a>"
                else:
                    para_content += f"<span style='{run_style}'>{run['text']}</span>"
            
            # Handle titles
            if paragraph.get('is_title'):
                text_content += f"<h1 class='title' style='{para_style}'>{para_content}</h1>"
            else:
                text_content += f"<p class='paragraph' style='{para_style}'>{para_content}</p>"
        
        # Handle background color
        bg_color = f"background-color: {shape_data.get('fill_color', 'transparent')};"
        return f"<div class='absolute flex flex-col items-start justify-start' style='{common_style}{bg_color}'>{text_content}</div>"
    return ""

def pptx_to_html(pptx_file, output_file, ignore_images=False):
    logging.info(f"Opening PowerPoint file: {pptx_file}")
    prs = Presentation(pptx_file)

    output_dir = os.path.dirname(output_file)
    image_dir = os.path.join(output_dir, "images")
    os.makedirs(image_dir, exist_ok=True)

    html_content = """
    <html>
    <head>
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <script src="https://cdn.tailwindcss.com"></script>
        <style>
            body { margin: 0; padding: 0; }
            .slide-container { max-width: 1024px; margin: 0 auto; padding: 2rem; }
            .slide { position: relative; width: 100%; padding-top: 75%; margin-bottom: 2rem; }
            .slide-content { position: absolute; top: 0; left: 0; right: 0; bottom: 0; }
            .absolute { position: absolute; }
            .paragraph { margin-bottom: 0.5em; }
            .slide-content > div { max-height: none !important; overflow: visible !important; }
            .title { font-size: 1.15em; font-weight: bold; margin-bottom: 0.5em; }
            a { text-decoration: underline; color: blue; }
            a:hover { text-decoration: none; }
        </style>
    </head>
    <body class="bg-gray-100">
    <div class="slide-container">
    """

    logging.info("Extracting slide data...")
    slide_data_list = [extract_slide_data(slide) for slide in prs.slides]

    logging.info("Processing slides...")
    with multiprocessing.Pool() as pool:
        try:
            slide_contents = list(tqdm(
                pool.starmap(partial(process_slide, ignore_images=ignore_images), 
                             [(json.dumps(data), image_dir, idx) for idx, data in enumerate(slide_data_list)]),
                total=len(slide_data_list),
                desc="Processing slides"
            ))
        except Exception as e:
            logging.error(f"An error occurred during slide processing: {str(e)}")
            slide_contents = []

    html_content += "".join(f'<div class="slide bg-white shadow-lg rounded-lg"><div class="slide-content">{content}</div></div>' for content in slide_contents)
    html_content += "</div></body></html>"

    logging.info(f"Writing HTML content to {output_file}")
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(html_content)

def main():
    parser = argparse.ArgumentParser(description="Convert PowerPoint (PPTX) to HTML")
    parser.add_argument("input", help="Input PPTX file")
    parser.add_argument("output", help="Output HTML file")
    parser.add_argument("-v", "--verbose", action="store_true", help="Enable verbose logging")
    parser.add_argument("-i", "--ignore-images", action="store_true", help="Ignore image export (for development)")
    args = parser.parse_args()

    log_level = logging.DEBUG if args.verbose else logging.INFO
    logging.basicConfig(level=log_level, format='%(asctime)s - %(levelname)s - %(message)s')

    logging.info("Starting PPTX to HTML conversion")
    pptx_to_html(args.input, args.output, ignore_images=args.ignore_images)
    logging.info(f"Conversion complete. HTML file saved as {args.output}")

if __name__ == "__main__":
    main()